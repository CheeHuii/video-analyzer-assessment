import json
from pathlib import Path
from typing import List, Dict, Any, Optional
from reportlab.lib.pagesizes import A4, landscape
from reportlab.lib.units import mm
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, PageBreak, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from PIL import Image as PILImage
import textwrap
import math
import os

# ---------- Helpers ----------
def load_transcript(transcript_path: str) -> Dict[str, Any]:
    with open(transcript_path, "r", encoding="utf-8") as f:
        return json.load(f)

def top_keypoints_from_transcript(transcript: Dict[str, Any], max_points: int = 6) -> List[str]:
    """
    Simple extractive algorithm:
     - Score segments by (confidence * length)
     - Return top N segment texts as 'key points'
    """
    segments = transcript.get("segments", [])
    scored = []
    for s in segments:
        conf = float(s.get("confidence", 1.0))
        length = max(1, len(s.get("text", "").split()))
        score = conf * math.log(1 + length)
        scored.append((score, s))
    scored.sort(key=lambda x: x[0], reverse=True)
    points = []
    for _, s in scored[:max_points]:
        txt = s.get("text", "").strip()
        if txt:
            points.append(txt)
    return points

def format_timestamp(sec: float) -> str:
    h = int(sec // 3600)
    m = int((sec % 3600) // 60)
    s = int(sec % 60)
    if h > 0:
        return f"{h:02d}:{m:02d}:{s:02d}"
    return f"{m:02d}:{s:02d}"

# ---------- PDF Generation ----------
def generate_pdf(video_dir: str, transcript_path: str, output_pdf: str, max_excerpt_chars: int = 300):
    """
    Create a multi-page PDF including:
      - Title page with metadata
      - Key points (bulleted)
      - Transcript excerpts with timestamps
      - Embedded graph images (if available)
    """
    video_dir = Path(video_dir)
    transcript = load_transcript(transcript_path)
    meta_path = video_dir / "meta.json"
    meta = {}
    if meta_path.exists():
        meta = json.loads(meta_path.read_text(encoding="utf-8"))

    doc = SimpleDocTemplate(output_pdf, pagesize=A4, rightMargin=20*mm, leftMargin=20*mm, topMargin=20*mm, bottomMargin=20*mm)
    styles = getSampleStyleSheet()
    story = []

    # Title page
    title_style = ParagraphStyle("Title", parent=styles["Title"], fontSize=24, alignment=1)
    story.append(Paragraph(meta.get("src_filename", "Video Summary"), title_style))
    story.append(Spacer(1, 6*mm))
    story.append(Paragraph(f"Video ID: {meta.get('video_id','-')}", styles["Normal"]))
    story.append(Paragraph(f"Duration: {meta.get('duration_s', '-') } seconds", styles["Normal"]))
    story.append(Paragraph(f"Ingested: {meta.get('ingested_at','-')}", styles["Normal"]))
    story.append(PageBreak())

    # Key points
    story.append(Paragraph("Key Points", styles["Heading2"]))
    points = top_keypoints_from_transcript(transcript, max_points=6)
    for p in points:
        story.append(Paragraph(f"• {p}", styles["BodyText"]))
        story.append(Spacer(1, 2*mm))
    story.append(PageBreak())

    # Transcript excerpts with timestamps (choose top few longest segments)
    story.append(Paragraph("Transcript Excerpts", styles["Heading2"]))
    segs = transcript.get("segments", [])
    # choose some excerpts: longest segments or highest score
    excerpts = sorted(segs, key=lambda s: (len(s.get("text","")), float(s.get("confidence",1.0))), reverse=True)[:10]
    for s in excerpts:
        start = format_timestamp(float(s.get("start",0)))
        end = format_timestamp(float(s.get("end",0)))
        header = Paragraph(f"{start} - {end}   (confidence: {s.get('confidence',1.0):.2f})", styles["Italic"])
        story.append(header)
        text = s.get("text","")
        # wrap long
        wrapped = "<br/>".join(textwrap.wrap(text, 120))
        story.append(Paragraph(wrapped, styles["BodyText"]))
        story.append(Spacer(1,3*mm))
    story.append(PageBreak())

    # Insert graphs/images (if any)
    graphs_dir = video_dir / "graphs"
    if graphs_dir.exists():
        story.append(Paragraph("Extracted Graphs", styles["Heading2"]))
        for img_path in sorted(graphs_dir.glob("*")):
            try:
                # Fit image to page width
                pil = PILImage.open(img_path)
                w, h = pil.size
                max_width = doc.width
                ratio = min(1.0, max_width / w)
                display_w = w * ratio
                display_h = h * ratio
                story.append(Image(str(img_path), width=display_w, height=display_h))
                story.append(Spacer(1,3*mm))
            except Exception as e:
                # skip unreadable images
                continue
    else:
        story.append(Paragraph("No extracted graphs available.", styles["Normal"]))

    # Build
    doc.build(story)
    return output_pdf

# ---------- PPTX Generation ----------
def add_slide_title(prs: Presentation, title_text: str, subtitle_text: Optional[str] = None):
    slide_layout = prs.slide_layouts[0]  # title slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text
    if subtitle_text:
        try:
            subtitle = slide.placeholders[1]
            subtitle.text = subtitle_text
        except Exception:
            pass
    return slide

def add_bullet_slide(prs: Presentation, title: str, bullets: List[str]):
    slide_layout = prs.slide_layouts[1]  # title and content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for b in bullets:
        p = body.add_paragraph()
        p.text = b
        p.level = 0
    return slide

def add_image_slide(prs: Presentation, title: str, img_path: str, caption: Optional[str] = None):
    """
    Adds an image slide safely; skips invalid or unreadable image files.
    """
    slide_layout = prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # optional title
    try:
        slide.shapes.title.text = title
    except Exception:
        pass

    # Validate image before adding
    try:
        with PILImage.open(img_path) as im:
            im.verify()  # check format validity
    except Exception as e:
        print(f"[WARN] Skipping invalid image: {img_path} ({e})")
        return slide  # skip this image silently

    left, top = Inches(1), Inches(1.5)
    try:
        slide.shapes.add_picture(img_path, left, top, width=Inches(8))
    except Exception as e:
        print(f"[WARN] Could not add image {img_path}: {e}")
        return slide

    if caption:
        tx = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(1))
        tf = tx.text_frame
        p = tf.add_paragraph()
        p.text = caption
        p.font.size = Pt(12)

    return slide

def generate_pptx(video_dir: str, transcript_path: str, output_pptx: str, max_points: int = 6):
    """
    Generate a PPTX with:
      - Title slide
      - Key points bullets
      - Slides per selected transcript excerpt
      - Slides for each graph
    """
    prs = Presentation()
    # set slide size optional: prs.slide_width, prs.slide_height

    video_dir = Path(video_dir)
    transcript = load_transcript(transcript_path)
    meta_path = video_dir / "meta.json"
    meta = {}
    if meta_path.exists():
        meta = json.loads(meta_path.read_text(encoding="utf-8"))

    # Title
    add_slide_title(prs, meta.get("src_filename", "Video Presentation"), f"Video ID: {meta.get('video_id','-')}")

    # Key points
    points = top_keypoints_from_transcript(transcript, max_points=max_points)
    add_bullet_slide(prs, "Key Points", points)

    # Excerpts -> slides (one excerpt per slide)
    segs = transcript.get("segments", [])
    excerpts = sorted(segs, key=lambda s: (len(s.get("text","")), float(s.get("confidence",1.0))), reverse=True)[:6]
    for s in excerpts:
        start = format_timestamp(float(s.get("start",0)))
        end = format_timestamp(float(s.get("end",0)))
        title = f"{start} - {end}"
        text = s.get("text","")
        add_bullet_slide(prs, title, textwrap.wrap(text, 80))

    # Graph slides
    graphs_dir = video_dir / "graphs"
    if graphs_dir.exists():
        graph_files = sorted([p for p in graphs_dir.glob("*") if p.is_file()])
        for gf in graph_files:
            add_image_slide(prs, "Graph / Chart", str(gf), caption=os.path.basename(gf))
    # Save
    prs.save(output_pptx)
    return output_pptx
